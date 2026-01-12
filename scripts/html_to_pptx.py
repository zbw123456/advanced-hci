#!/usr/bin/env python3
import sys
from html.parser import HTMLParser
from pptx import Presentation
from pptx.util import Inches, Pt

class SlideHTMLParser(HTMLParser):
    def __init__(self):
        super().__init__()
        self.in_section = False
        self.slides = []
        self.current = []
        self.last_tag = None
        self.ignoring = False

    def handle_starttag(self, tag, attrs):
        self.last_tag = tag
        if tag == 'section':
            self.in_section = True
            self.current = []
        if tag in ('script', 'style'):
            self.ignoring = True

    def handle_endtag(self, tag):
        self.last_tag = None
        if tag == 'section':
            self.in_section = False
            text = '\n'.join(self.current).strip()
            if text:
                self.slides.append(text)
        if tag in ('script', 'style'):
            self.ignoring = False

    def handle_data(self, data):
        if self.ignoring:
            return
        if not self.in_section:
            return
        s = data.strip()
        if not s:
            return
        # Insert a line break before common block tags so items become separate lines
        if self.last_tag in ('li', 'p', 'div') and self.current:
            self.current.append('\n')
        self.current.append(s)

def html_to_pptx(html_path, out_path):
    parser = SlideHTMLParser()
    with open(html_path, 'r', encoding='utf-8') as f:
        parser.feed(f.read())

    if not parser.slides:
        print('No slides found in', html_path)
        return 1

    prs = Presentation()
    for slide_text in parser.slides:
        # normalize and split lines
        parts = [line.strip() for line in slide_text.split('\n') if line.strip()]
        if not parts:
            continue

        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        left = Inches(0.6)
        top = Inches(0.5)
        width = Inches(9)
        title_height = Inches(1.2)

        title_box = slide.shapes.add_textbox(left, top, width, title_height)
        title_tf = title_box.text_frame
        title_tf.clear()
        title_p = title_tf.paragraphs[0]
        title_p.text = parts[0]
        title_p.font.size = Pt(28)
        title_p.font.bold = True

        # content
        content_top = top + title_height + Inches(0.3)
        content_box = slide.shapes.add_textbox(left, content_top, width, Inches(5))
        content_tf = content_box.text_frame
        content_tf.clear()

        if len(parts) >= 2:
            # if second line is short, treat as subtitle then bullets
            subtitle = parts[1]
            if len(subtitle) < 140 and len(parts) > 2:
                p = content_tf.paragraphs[0]
                p.text = subtitle
                p.font.size = Pt(16)
                # add bullets
                for b in parts[2:]:
                    para = content_tf.add_paragraph()
                    para.text = b
                    para.level = 0
                    para.font.size = Pt(14)
            else:
                # treat remaining as paragraphs/bullets
                for b in parts[1:]:
                    para = content_tf.paragraphs[0] if content_tf.paragraphs[0].text == '' else content_tf.add_paragraph()
                    para.text = b
                    para.level = 0
                    para.font.size = Pt(14)

    prs.save(out_path)
    print('Saved PPTX to', out_path)
    return 0

if __name__ == '__main__':
    html_path = sys.argv[1] if len(sys.argv) > 1 else 'literature_review.html'
    out_path = sys.argv[2] if len(sys.argv) > 2 else 'literature_review.pptx'
    sys.exit(html_to_pptx(html_path, out_path))
