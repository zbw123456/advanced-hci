
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()

    # Define a function to add a slide with title and content
    def add_slide(title, content_lines, notes=""):
        # Layout 1 is usually Title + Content
        slide_layout = prs.slide_layouts[1] 
        slide = prs.slides.add_slide(slide_layout)
        
        # Set Title
        title_shape = slide.shapes.title
        title_shape.text = title
        
        # Set Content
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.word_wrap = True
        
        # Clear existing paragraphs
        for p in tf.paragraphs:
            p.text = ""
            
        # Add content
        for i, line in enumerate(content_lines):
            if i == 0 and not tf.paragraphs[0].text:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            # Check for sub-bullets
            if line.startswith("    "):
                p.text = line.strip()
                p.level = 1
            else:
                p.text = line.strip()
                p.level = 0
                
            p.font.size = Pt(20)

        # Add notes
        if notes:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = notes
    
    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "MindCare: A Privacy-First Multimodal Emotion Monitoring System"
    subtitle.text = "Advanced HCI Final Project\nv2.0 Update"

    # Slide 2: Field & Problem
    add_slide(
        "Field & Problem Identification",
        [
            "The Field: Affective Computing & HCI applied to Mental Health.",
            "The Problem: 'Invisible Stress'",
            "    Low Self-Awareness: Missing early signs of burnout.",
            "    Lack of Tools:",
            "        Wearables are intrusive/expensive.",
            "        Manual journaling suffers from bias.",
            "    Privacy Concerns: Distrust of cloud-based analysis.",
            "Goal: A locally processing desktop companion."
        ],
        notes="Focus on the gap between expensive wearables and manual apps."
    )

    # Slide 3: Literature Review 1
    add_slide(
        "Literature Review: Current Solutions",
        [
            "1. Physiological Sensors (Wearables)",
            "    Pros: High accuracy (HRV, GSR).",
            "    Cons: High barrier, charging fatigue, discomfort.",
            "2. Self-Reporting (Apps)",
            "    Pros: Captures subjective experience.",
            "    Cons: High friction, data gaps, retrospective bias."
        ]
    )

    # Slide 4: Literature Review 2
    add_slide(
        "Literature Review: Technical Approaches",
        [
            "Facial Expression Recognition (FER)",
            "    Traditional ML: Haar Cascades + SVM. Fast but less accurate.",
            "    Deep Learning: CNNs (VGG-16, ResNet). High accuracy but heavy.",
            "The Gap:",
            "    Most SOTA models require GPUs or Cloud APIs.",
            "    Need for lightweight, privacy-preserving local implementation."
        ]
    )

    # Slide 5: Solution Overview
    add_slide(
        "Our Solution: MindCare",
        [
            "Concept: Non-intrusive desktop app monitoring emotion via webcam.",
            "The ARCADE Framework:",
            "    Awareness: Real-time visual feedback.",
            "    Responsiveness: Immediate alerts.",
            "    Control: Local encrypted storage.",
            "Key Innovation:",
            "    'Privacy by Design' - No video frames leave the device."
        ]
    )

    # Slide 6: Functionality - Real-Time
    add_slide(
        "Functionality: Real-Time Monitoring",
        [
            "Vision Pipeline:",
            "    Face Detection (Haar Cascades).",
            "    Smile Detection (Geometric analysis).",
            "Visual Feedback:",
            "    'Searching...' (Red status).",
            "    'Face Detected' (Green status).",
            "    'Smiling! :)' (Yellow reactive status)."
        ]
    )

    # Slide 7: Functionality - Dashboard
    add_slide(
        "Functionality: Interactive Dashboard",
        [
            "Built with PyQt5 for Professional UX.",
            "Dynamic Metrics:",
            "    Emotion Probabilities (Bar Charts).",
            "    Stress Index (Calculated metric).",
            "    Trend Buffer (15-min sliding window)."
        ]
    )

    # Slide 8: Functionality - Security
    add_slide(
        "Functionality: Security & Privacy",
        [
            "Local Processing Only (No Cloud).",
            "Encryption:",
            "    Uses AES-128 via cryptography.fernet.",
            "    Session logs encrypted at rest.",
            "Key Management:",
            "    Keys generated and stored locally."
        ]
    )
    
    # Slide 9: Functionality - FSM
    add_slide(
        "Functionality: System Control",
        [
            "Finite State Machine (FSM) driven.",
            "States:",
            "    Initializing",
            "    Monitoring (Active Loop)",
            "    Paused (Privacy Mode)",
            "    Alerting (Stress > Threshold)",
            "Error Recovery: Handles camera disconnects gracefully."
        ]
    )

    # Slide 10: Persona
    add_slide(
        "User Persona: Maria",
        [
            "Profile: Maria Chen (28), Senior Developer.",
            "Behaviors:",
            "    8+ hours screen time.",
            "    Enters 'flow' but forgets self-care.",
            "    Privacy conscious (covers webcam).",
            "Pain Points:",
            "    Unexpected burnout.",
            "    Physical tension headaches."
        ]
    )

    # Slide 11: Scenario 1
    add_slide(
        "Scenario: Morning Routine",
        [
            "09:00 AM: Start",
            "    Maria launches MindCare background process.",
            "    Green light confirms system ready.",
            "11:30 AM: Deep Work",
            "    Maria encounters a bug.",
            "    Unconsciously frowns/tenses for 45 mins."
        ]
    )
    
    # Slide 12: Scenario 2
    add_slide(
        "Scenario: Intervention",
        [
            "11:45 AM: Detection",
            "    System detects negative valence > 0.8.",
            "    Triggers Stress Alert.",
            "11:50 AM: Action",
            "    Notification: 'High stress detected. Stretch?'",
            "    Maria takes a break, smiles at a meme.",
            "    System logs 'Happy' spike, stress resets."
        ]
    )

    # Slide 13: Hardware Components
    add_slide(
        "Hardware Requirements",
        [
            "Minimalist Setup (No Wearables)",
            "Camera:",
            "    Standard Webcam (720p, 30fps).",
            "Computing:",
            "    Dual Core CPU (2.0GHz+).",
            "    8GB RAM.",
            "Storage:",
            "    <100MB for application + logs."
        ]
    )

    # Slide 14: Technical Tools 1
    add_slide(
        "Technical Stack",
        [
            "Language: Python 3.9+",
            "GUI: PyQt5 (Signals/Slots architecture)",
            "Computer Vision: OpenCV (cv2)",
            "Security: Cryptography (Fernet AES)",
            "Data: NumPy, SQLite (JSON logs)"
        ]
    )

    # Slide 15: Technical Tools 2 (Architecture)
    add_slide(
        "Architecture Pattern",
        [
            "Model-View-Controller (MVC)",
            "    Model: FSM Logic, Security Manager",
            "    View: MainWindow, VideoWidget (PyQt5)",
            "    Controller: MainApp, CameraThread",
            "Threading:",
            "    QThread used for non-blocking UI.",
            "    Signals for thread-safe updates."
        ]
    )

    # Slide 16: UML Component Diagram
    add_slide(
        "UML Component Diagram",
        [
            "[ Diagram Description ]",
            "Presentation Layer (UI)",
            "    |",
            "    v",
            "Application Layer (Main Control)",
            "    |",
            "    v",
            "Processing Layer (CV/ML)",
            "    |",
            "    v",
            "Hardware Layer (Camera)"
        ]
    )

    # Slide 17: UML State Chart
    add_slide(
        "UML State Chart Diagram",
        [
            "States Flow:",
            "    [Start] -> [Initializing] -> [Monitoring]",
            "    [Monitoring] <-> [Paused]",
            "    [Monitoring] -> [Alerting] (Stress > 0.8)",
            "    [Monitoring] -> [Stopped] -> [Exit]"
        ]
    )

    # Slide 18: Implementation Results
    add_slide(
        "Implementation: Performance",
        [
            "Latency: < 50ms end-to-end.",
            "Responsiveness: UI remains fluid (60fps) via Threading.",
            "Accuracy:",
            "    Face Detection: >95% (frontal).",
            "    Smile Detection: Reliable trigger for 'Happy' state."
        ]
    )
    
    # Slide 19: Implementation - Visuals
    add_slide(
        "Implementation: Visual Interface",
        [
            "Dark Mode Design for reduced eye strain.",
            "Split View:",
            "    Left: Real-time Camera Feed + Overlays.",
            "    Right: Analytical Dashboards.",
            "Key Feedback:",
            "    'Searching...' text prevents user confusion."
        ]
    )

    # Slide 20: Conclusions & Lessons
    add_slide(
        "Conclusions & Lessons Learnt",
        [
            "Conclusions:",
            "    Privacy-preserving desktop monitoring is viable.",
            "    Reactive algorithms effectively simulate awareness.",
            "Lessons:",
            "    Concurrency is critical for GUI+CV.",
            "    Clear visual feedback is essential for UX."
        ]
    )
    
    # Slide 21: References
    add_slide(
        "References",
        [
            "1. Ekman, P. (1992). An argument for basic emotions.",
            "2. Calvo, R. A. (2010). Affect detection review.",
            "3. Viola, P., & Jones, M. (2001). Rapid object detection (CVPR).",
            "4. PyQt5 Documentation.",
            "5. OpenAI / OpenCV Community Resources."
        ]
    )

    output_file = "MindCare_Final_Presentation_v2.pptx"
    prs.save(output_file)
    print(f"Successfully generated {output_file}")

if __name__ == "__main__":
    create_presentation()
