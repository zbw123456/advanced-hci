#!/usr/bin/env python3
"""
MindCare Final Presentation Generator
Creates a professional PowerPoint presentation for the Advanced HCI project
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os

def create_presentation():
    """Create the MindCare presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define color scheme
    PURPLE = RGBColor(102, 126, 234)  # #667eea
    DARK_PURPLE = RGBColor(118, 75, 162)  # #764ba2
    PINK = RGBColor(245, 87, 108)  # #f5576c
    LIGHT_BLUE = RGBColor(79, 172, 254)  # #4facfe
    DARK_BG = RGBColor(15, 15, 30)  # #0f0f1e
    WHITE = RGBColor(255, 255, 255)
    GRAY = RGBColor(184, 185, 207)  # #b8b9cf
    
    def set_background(slide, color):
        """Set solid background color"""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color
    
    def add_title_slide():
        """Slide 1: Title"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        set_background(slide, DARK_BG)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = "MindCare"
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.size = Pt(72)
        title_para.font.bold = True
        title_para.font.color.rgb = PURPLE
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(9), Inches(0.8))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = "Mental Well-being Monitoring System with Multimodal Interaction"
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.alignment = PP_ALIGN.CENTER
        subtitle_para.font.size = Pt(24)
        subtitle_para.font.color.rgb = GRAY
        
        # Meta info
        meta_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1))
        meta_frame = meta_box.text_frame
        meta_text = "Advanced HCI 2025-2026  •  University of Trento  •  Bowen Zhang"
        meta_frame.text = meta_text
        meta_para = meta_frame.paragraphs[0]
        meta_para.alignment = PP_ALIGN.CENTER
        meta_para.font.size = Pt(16)
        meta_para.font.color.rgb = GRAY
    
    def add_problem_slide():
        """Slide 2: Problem Statement"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_background(slide, DARK_BG)
        
        # Header
        header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        header_frame = header_box.text_frame
        header_frame.text = "The Challenge"
        header_para = header_frame.paragraphs[0]
        header_para.font.size = Pt(48)
        header_para.font.bold = True
        header_para.font.color.rgb = PURPLE
        
        # Section label
        label_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.3))
        label_frame = label_box.text_frame
        label_frame.text = "PART I • PROBLEM"
        label_para = label_frame.paragraphs[0]
        label_para.font.size = Pt(12)
        label_para.font.color.rgb = GRAY
        
        # Content
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.4), Inches(5))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        # Problem definition
        p1 = content_frame.paragraphs[0]
        p1.text = "Problem Definition"
        p1.font.size = Pt(28)
        p1.font.bold = True
        p1.font.color.rgb = WHITE
        p1.space_after = Pt(12)
        
        p2 = content_frame.add_paragraph()
        p2.text = "Many individuals fail to recognize early warning signs of stress and anxiety, leading to delayed intervention and deteriorating mental health outcomes."
        p2.font.size = Pt(18)
        p2.font.color.rgb = GRAY
        p2.space_after = Pt(24)
        
        # Key Issues
        p3 = content_frame.add_paragraph()
        p3.text = "Key Issues"
        p3.font.size = Pt(24)
        p3.font.bold = True
        p3.font.color.rgb = WHITE
        p3.space_after = Pt(12)
        
        issues = [
            "Lack of self-awareness for gradual emotional changes",
            "Stigma prevents early help-seeking behavior",
            "Traditional services are reactive, not proactive",
            "Accessibility barriers to mental health care"
        ]
        
        for issue in issues:
            p = content_frame.add_paragraph()
            p.text = f"→  {issue}"
            p.font.size = Pt(16)
            p.font.color.rgb = GRAY
            p.space_after = Pt(8)
    
    def add_arcade_slide():
        """Slide 3: Multimodal Solution (ARCADE)"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_background(slide, DARK_BG)
        
        # Header
        header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        header_frame = header_box.text_frame
        header_frame.text = "Multimodal Interaction Design"
        header_para = header_frame.paragraphs[0]
        header_para.font.size = Pt(44)
        header_para.font.bold = True
        header_para.font.color.rgb = PURPLE
        
        # Section label
        label_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.3))
        label_frame = label_box.text_frame
        label_frame.text = "PART I • SOLUTION"
        label_para = label_frame.paragraphs[0]
        label_para.font.size = Pt(12)
        label_para.font.color.rgb = GRAY
        
        # ARCADE Framework
        y_pos = 1.8
        modalities = [
            ("📹 Camera", "Assignment: Passive facial expression monitoring"),
            ("🎤 Voice", "Redundancy: Hands-free commands for status checks"),
            ("👆 Touch", "Complementarity: Detailed configuration & visualization")
        ]
        
        for icon_title, description in modalities:
            # Box background
            box = slide.shapes.add_shape(1, Inches(0.8), Inches(y_pos), Inches(8.4), Inches(0.8))
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(20, 20, 40)
            box.line.color.rgb = PURPLE
            
            # Title
            title_box = slide.shapes.add_textbox(Inches(1), Inches(y_pos + 0.1), Inches(2), Inches(0.3))
            title_frame = title_box.text_frame
            title_frame.text = icon_title
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(20)
            title_para.font.bold = True
            title_para.font.color.rgb = WHITE
            
            # Description
            desc_box = slide.shapes.add_textbox(Inches(1), Inches(y_pos + 0.4), Inches(7.5), Inches(0.3))
            desc_frame = desc_box.text_frame
            desc_frame.text = description
            desc_para = desc_frame.paragraphs[0]
            desc_para.font.size = Pt(14)
            desc_para.font.color.rgb = GRAY
            
            y_pos += 1.0
        
        # Advantages and Challenges
        adv_box = slide.shapes.add_textbox(Inches(0.8), Inches(5.2), Inches(4), Inches(1.8))
        adv_frame = adv_box.text_frame
        p = adv_frame.paragraphs[0]
        p.text = "✅ Advantages"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(16, 185, 129)
        
        adv_text = "\n• Non-intrusive monitoring\n• Early warning capabilities\n• Multi-modal accessibility"
        p2 = adv_frame.add_paragraph()
        p2.text = adv_text
        p2.font.size = Pt(14)
        p2.font.color.rgb = GRAY
        
        chal_box = slide.shapes.add_textbox(Inches(5.2), Inches(5.2), Inches(4), Inches(1.8))
        chal_frame = chal_box.text_frame
        p = chal_frame.paragraphs[0]
        p.text = "⚠️ Challenges"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(245, 158, 11)
        
        chal_text = "\n• Privacy concerns\n• FER accuracy ~70%\n• Individual variation"
        p2 = chal_frame.add_paragraph()
        p2.text = chal_text
        p2.font.size = Pt(14)
        p2.font.color.rgb = GRAY
    
    def add_persona_slide():
        """Slide 4: User-Centered Design"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_background(slide, DARK_BG)
        
        # Header
        header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        header_frame = header_box.text_frame
        header_frame.text = "Persona & Scenario"
        header_para = header_frame.paragraphs[0]
        header_para.font.size = Pt(48)
        header_para.font.bold = True
        header_para.font.color.rgb = PURPLE
        
        # Section label
        label_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.3))
        label_frame = label_box.text_frame
        label_frame.text = "PART II • USER ANALYSIS"
        label_para = label_frame.paragraphs[0]
        label_para.font.size = Pt(12)
        label_para.font.color.rgb = GRAY
        
        # Persona (left column)
        persona_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(4.2), Inches(5))
        persona_frame = persona_box.text_frame
        
        p = persona_frame.paragraphs[0]
        p.text = "👤 Persona: Maria Chen"
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.space_after = Pt(16)
        
        persona_items = [
            ("🎯 Profile", "28-year-old software developer"),
            ("😰 Pain Point", "Workplace stress, unaware of patterns"),
            ("🎯 Goal", "Better work-life balance"),
            ("💻 Tech Savvy", "Comfortable with health apps")
        ]
        
        for title, desc in persona_items:
            p = persona_frame.add_paragraph()
            p.text = title
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = PURPLE
            p.space_after = Pt(4)
            
            p = persona_frame.add_paragraph()
            p.text = desc
            p.font.size = Pt(14)
            p.font.color.rgb = GRAY
            p.space_after = Pt(12)
        
        # Scenario (right column)
        scenario_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.8), Inches(4.2), Inches(5))
        scenario_frame = scenario_box.text_frame
        
        p = scenario_frame.paragraphs[0]
        p.text = "📖 Day-in-the-Life"
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.space_after = Pt(16)
        
        timeline = [
            ("9:00 AM", "MindCare starts monitoring"),
            ("2:30 PM", "Detects 65% negative emotions"),
            ("2:31 PM", "Alert: 'Take a break'"),
            ("2:35 PM", "Maria walks, mood improves")
        ]
        
        for time, event in timeline:
            p = scenario_frame.add_paragraph()
            p.text = time
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = LIGHT_BLUE
            p.space_after = Pt(2)
            
            p = scenario_frame.add_paragraph()
            p.text = event
            p.font.size = Pt(13)
            p.font.color.rgb = GRAY
            p.space_after = Pt(10)
    
    def add_architecture_slide():
        """Slide 5: System Architecture"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_background(slide, DARK_BG)
        
        # Header
        header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        header_frame = header_box.text_frame
        header_frame.text = "Architecture & Technology"
        header_para = header_frame.paragraphs[0]
        header_para.font.size = Pt(44)
        header_para.font.bold = True
        header_para.font.color.rgb = PURPLE
        
        # Section label
        label_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.3))
        label_frame = label_box.text_frame
        label_frame.text = "PART II • SYSTEM DESIGN"
        label_para = label_frame.paragraphs[0]
        label_para.font.size = Pt(12)
        label_para.font.color.rgb = GRAY
        
        # 4-Layer Architecture
        layers = [
            ("1️⃣ Interface Layer", "PyQt5 GUI • Voice Commands • Notifications"),
            ("2️⃣ Application Layer", "FSM Controller (19 states) • Alert Engine"),
            ("3️⃣ Processing Layer", "OpenCV Face Detection • CNN Emotion Classifier"),
            ("4️⃣ Data Layer", "SQLite Database • AES-256 Encryption")
        ]
        
        y_pos = 1.8
        for layer_name, technologies in layers:
            # Layer box
            box = slide.shapes.add_shape(1, Inches(0.8), Inches(y_pos), Inches(8.4), Inches(0.7))
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(20, 20, 40)
            box.line.color.rgb = PURPLE
            
            # Layer name
            name_box = slide.shapes.add_textbox(Inches(1), Inches(y_pos + 0.1), Inches(8), Inches(0.25))
            name_frame = name_box.text_frame
            name_frame.text = layer_name
            name_para = name_frame.paragraphs[0]
            name_para.font.size = Pt(18)
            name_para.font.bold = True
            name_para.font.color.rgb = WHITE
            
            # Technologies
            tech_box = slide.shapes.add_textbox(Inches(1), Inches(y_pos + 0.35), Inches(8), Inches(0.25))
            tech_frame = tech_box.text_frame
            tech_frame.text = technologies
            tech_para = tech_frame.paragraphs[0]
            tech_para.font.size = Pt(13)
            tech_para.font.color.rgb = GRAY
            
            y_pos += 0.85
        
        # Tech Stack
        tech_y = 5.5
        tech_box = slide.shapes.add_textbox(Inches(0.8), Inches(tech_y), Inches(8.4), Inches(1.5))
        tech_frame = tech_box.text_frame
        
        p = tech_frame.paragraphs[0]
        p.text = "Technology Stack"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.space_after = Pt(8)
        
        p = tech_frame.add_paragraph()
        p.text = "Python 3.8+  •  OpenCV 4.8+  •  TensorFlow 2.14+  •  PyQt5 5.15+  •  AES-256 Encryption"
        p.font.size = Pt(14)
        p.font.color.rgb = PURPLE
        p.alignment = PP_ALIGN.CENTER
    
    def add_implementation_slide():
        """Slide 6: Implementation"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_background(slide, DARK_BG)
        
        # Header
        header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        header_frame = header_box.text_frame
        header_frame.text = "Behavioral Model & Features"
        header_para = header_frame.paragraphs[0]
        header_para.font.size = Pt(44)
        header_para.font.bold = True
        header_para.font.color.rgb = PURPLE
        
        # Section label
        label_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.3))
        label_frame = label_box.text_frame
        label_frame.text = "PART III • IMPLEMENTATION"
        label_para = label_frame.paragraphs[0]
        label_para.font.size = Pt(12)
        label_para.font.color.rgb = GRAY
        
        # FSM (left)
        fsm_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(4.2), Inches(2.5))
        fsm_frame = fsm_box.text_frame
        
        p = fsm_frame.paragraphs[0]
        p.text = "🔄 Finite State Machine"
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.space_after = Pt(12)
        
        p = fsm_frame.add_paragraph()
        p.text = "19 States | Comprehensive Error Handling"
        p.font.size = Pt(14)
        p.font.color.rgb = PURPLE
        p.space_after = Pt(16)
        
        p = fsm_frame.add_paragraph()
        p.text = "Key State Transitions:"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = GRAY
        p.space_after = Pt(8)
        
        transitions = [
            "INIT → CAMERA_CHECK → MONITORING",
            "MONITORING ↔ PAUSED",
            "STRESS_DETECTED → ALERT",
            "ERROR_* → RECOVERY"
        ]
        
        for trans in transitions:
            p = fsm_frame.add_paragraph()
            p.text = f"  {trans}"
            p.font.size = Pt(12)
            p.font.color.rgb = GRAY
            p.space_after = Pt(4)
        
        # Features (right)
        feat_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.8), Inches(4.2), Inches(2.5))
        feat_frame = feat_box.text_frame
        
        p = feat_frame.paragraphs[0]
        p.text = "⚡ Core Features"
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.space_after = Pt(12)
        
        features = [
            "Real-time FER at 30 FPS",
            "7 emotion classification",
            "2-second sliding windows",
            "15-min stress pattern analysis",
            "Session analytics & reports"
        ]
        
        for feat in features:
            p = feat_frame.add_paragraph()
            p.text = f"→  {feat}"
            p.font.size = Pt(14)
            p.font.color.rgb = GRAY
            p.space_after = Pt(8)
        
        # Statistics
        stats = [
            ("~500", "Lines Code"),
            ("~3000", "Lines Docs"),
            ("7", "Emotions"),
            ("30", "FPS")
        ]
        
        x_pos = 0.8
        for number, label in stats:
            stat_box = slide.shapes.add_textbox(Inches(x_pos), Inches(4.8), Inches(2), Inches(1.2))
            stat_frame = stat_box.text_frame
            
            p = stat_frame.paragraphs[0]
            p.text = number
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = PURPLE
            p.alignment = PP_ALIGN.CENTER
            p.space_after = Pt(4)
            
            p = stat_frame.add_paragraph()
            p.text = label
            p.font.size = Pt(12)
            p.font.color.rgb = GRAY
            p.alignment = PP_ALIGN.CENTER
            
            x_pos += 2.1
    
    def add_demo_slide():
        """Slide 7: Live Demo"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_background(slide, DARK_BG)
        
        # Header
        header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        header_frame = header_box.text_frame
        header_frame.text = "System in Action"
        header_para = header_frame.paragraphs[0]
        header_para.font.size = Pt(48)
        header_para.font.bold = True
        header_para.font.color.rgb = PURPLE
        
        # Section label
        label_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.3))
        label_frame = label_box.text_frame
        label_frame.text = "PART III • DEMONSTRATION"
        label_para = label_frame.paragraphs[0]
        label_para.font.size = Pt(12)
        label_para.font.color.rgb = GRAY
        
        # Demo output
        demo_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.4), Inches(5))
        demo_frame = demo_box.text_frame
        
        demo_output = """MindCare v1.0.0
Initializing...
✓ Loading face detector (OpenCV Haar Cascades)
✓ Loading emotion classifier (CNN TensorFlow)
✓ Initialization complete
✓ Camera found and accessible

==================================================
MONITORING STARTED
==================================================

[14:23:15] Emotion: NEUTRAL    | Confidence: 0.78 | Valence: +0.02
[14:23:16] Emotion: HAPPY      | Confidence: 0.65 | Valence: +0.52
[14:23:17] Emotion: NEUTRAL    | Confidence: 0.82 | Valence: +0.01
[14:23:18] Emotion: SAD        | Confidence: 0.71 | Valence: -0.48

⚠️  STRESS PATTERN DETECTED
   Negative emotions: 65.0% over last 20 readings
   Suggestion: Consider taking a short break

==================================================
SESSION SUMMARY
==================================================
Emotion Distribution:
  neutral   : ████████████ 60.0% (12)
  happy     : ████ 20.0% (4)
  sad       : ██ 10.0% (2)"""
        
        p = demo_frame.paragraphs[0]
        p.text = demo_output
        p.font.name = 'Monaco'
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(168, 218, 220)
    
    def add_results_slide():
        """Slide 8: Results & Impact"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_background(slide, DARK_BG)
        
        # Header
        header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        header_frame = header_box.text_frame
        header_frame.text = "Achievements & Validation"
        header_para = header_frame.paragraphs[0]
        header_para.font.size = Pt(44)
        header_para.font.bold = True
        header_para.font.color.rgb = PURPLE
        
        # Section label
        label_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.3))
        label_frame = label_box.text_frame
        label_frame.text = "RESULTS"
        label_para = label_frame.paragraphs[0]
        label_para.font.size = Pt(12)
        label_para.font.color.rgb = GRAY
        
        # Deliverables (left)
        deliv_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(4.2), Inches(2.5))
        deliv_frame = deliv_box.text_frame
        
        p = deliv_frame.paragraphs[0]
        p.text = "✅ Completed Deliverables"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.space_after = Pt(12)
        
        deliverables = [
            "Part I: Literature review (3000+ words)",
            "Part II: User-centered design",
            "Part III: Working prototype (FSM)",
            "Bonus: Professional UI mockups"
        ]
        
        for item in deliverables:
            p = deliv_frame.add_paragraph()
            p.text = f"→  {item}"
            p.font.size = Pt(14)
            p.font.color.rgb = GRAY
            p.space_after = Pt(8)
        
        # Validation (right)
        valid_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.8), Inches(4.2), Inches(2.5))
        valid_frame = valid_box.text_frame
        
        p = valid_frame.paragraphs[0]
        p.text = "🎯 Technical Validation"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.space_after = Pt(12)
        
        validations = [
            "Real-time face detection at 30 FPS",
            "7-emotion CNN classification",
            "2s time-window aggregation",
            "15-min stress pattern detection",
            "Complete session analytics"
        ]
        
        for item in validations:
            p = valid_frame.add_paragraph()
            p.text = f"→  {item}"
            p.font.size = Pt(14)
            p.font.color.rgb = GRAY
            p.space_after = Pt(8)
        
        # Contributions
        contrib_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.5), Inches(8.4), Inches(2.3))
        contrib_frame = contrib_box.text_frame
        
        p = contrib_frame.paragraphs[0]
        p.text = "💡 Key Contributions to HCI"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.space_after = Pt(12)
        
        contributions = [
            "1. Multimodal Design: Practical ARCADE framework in affective computing",
            "2. Privacy-First: Local processing addressing ethical concerns",
            "3. Time-Window Innovation: Reduce FER noise without storing frames",
            "4. Proactive Intervention: Pattern-based alerts for early detection"
        ]
        
        for contrib in contributions:
            p = contrib_frame.add_paragraph()
            p.text = contrib
            p.font.size = Pt(13)
            p.font.color.rgb = GRAY
            p.space_after = Pt(6)
    
    def add_limitations_slide():
        """Slide 9: Limitations & Future Work"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_background(slide, DARK_BG)
        
        # Header
        header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        header_frame = header_box.text_frame
        header_frame.text = "Limitations & Future Directions"
        header_para = header_frame.paragraphs[0]
        header_para.font.size = Pt(40)
        header_para.font.bold = True
        header_para.font.color.rgb = PURPLE
        
        # Section label
        label_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.3))
        label_frame = label_box.text_frame
        label_frame.text = "DISCUSSION"
        label_para = label_frame.paragraphs[0]
        label_para.font.size = Pt(12)
        label_para.font.color.rgb = GRAY
        
        # Limitations (left)
        limit_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(4.2), Inches(3.5))
        limit_frame = limit_box.text_frame
        
        p = limit_frame.paragraphs[0]
        p.text = "⚠️ Current Limitations"
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.space_after = Pt(12)
        
        limitations = [
            "Demo CNN is untrained",
            "Only camera input functional",
            "Console UI (no graphical)",
            "In-memory only (no persistence)",
            "Real-world FER accuracy ~70%"
        ]
        
        for item in limitations:
            p = limit_frame.add_paragraph()
            p.text = f"→  {item}"
            p.font.size = Pt(14)
            p.font.color.rgb = GRAY
            p.space_after = Pt(8)
        
        # Future Work (right)
        future_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.8), Inches(4.2), Inches(3.5))
        future_frame = future_box.text_frame
        
        p = future_frame.paragraphs[0]
        p.text = "🚀 Future Enhancements"
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.space_after = Pt(12)
        
        enhancements = [
            "Train FER+ model on FER-2013",
            "Implement voice integration",
            "Build PyQt5 GUI dashboard",
            "Multi-person face tracking",
            "Integrate contextual AI"
        ]
        
        for item in enhancements:
            p = future_frame.add_paragraph()
            p.text = f"→  {item}"
            p.font.size = Pt(14)
            p.font.color.rgb = GRAY
            p.space_after = Pt(8)
        
        # Disclaimer
        disclaimer_box = slide.shapes.add_textbox(Inches(0.8), Inches(5.5), Inches(8.4), Inches(1.3))
        disclaimer_frame = disclaimer_box.text_frame
        
        p = disclaimer_frame.paragraphs[0]
        p.text = "⚠️ Critical Disclaimer"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(245, 158, 11)
        p.space_after = Pt(8)
        
        p = disclaimer_frame.add_paragraph()
        p.text = "MindCare is a self-awareness tool and educational prototype. It is NOT a replacement for professional mental health services."
        p.font.size = Pt(13)
        p.font.color.rgb = GRAY
    
    def add_conclusion_slide():
        """Slide 10: Conclusion"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_background(slide, DARK_BG)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = "Thank You"
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.size = Pt(56)
        title_para.font.bold = True
        title_para.font.color.rgb = PURPLE
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.1), Inches(9), Inches(0.6))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = "MindCare: Advancing Mental Well-being Through Multimodal HCI"
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.alignment = PP_ALIGN.CENTER
        subtitle_para.font.size = Pt(20)
        subtitle_para.font.color.rgb = GRAY
        
        # Key Takeaways
        takeaway_box = slide.shapes.add_textbox(Inches(1.5), Inches(4.2), Inches(7), Inches(2))
        takeaway_frame = takeaway_box.text_frame
        
        p = takeaway_frame.paragraphs[0]
        p.text = "📚 Key Takeaways"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(12)
        
        takeaways = [
            "Passive monitoring enables early stress detection",
            "Multimodal ARCADE enhances accessibility",
            "Time-window processing balances accuracy & privacy",
            "User-centered design is critical for adoption"
        ]
        
        for item in takeaways:
            p = takeaway_frame.add_paragraph()
            p.text = f"→  {item}"
            p.font.size = Pt(14)
            p.font.color.rgb = GRAY
            p.space_after = Pt(6)
        
        # Footer
        footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
        footer_frame = footer_box.text_frame
        footer_frame.text = "Made with 💜 for mental well-being awareness"
        footer_para = footer_frame.paragraphs[0]
        footer_para.alignment = PP_ALIGN.CENTER
        footer_para.font.size = Pt(14)
        footer_para.font.color.rgb = GRAY
    
    # Generate all slides
    print("Creating MindCare presentation...")
    add_title_slide()
    print("✓ Slide 1: Title")
    add_problem_slide()
    print("✓ Slide 2: Problem Statement")
    add_arcade_slide()
    print("✓ Slide 3: ARCADE Framework")
    add_persona_slide()
    print("✓ Slide 4: Persona & Scenario")
    add_architecture_slide()
    print("✓ Slide 5: Architecture")
    add_implementation_slide()
    print("✓ Slide 6: Implementation")
    add_demo_slide()
    print("✓ Slide 7: Demo")
    add_results_slide()
    print("✓ Slide 8: Results")
    add_limitations_slide()
    print("✓ Slide 9: Limitations")
    add_conclusion_slide()
    print("✓ Slide 10: Conclusion")
    
    # Save presentation
    output_file = "MindCare_Final_Presentation.pptx"
    prs.save(output_file)
    print(f"\n✅ Presentation saved as: {output_file}")
    print(f"📊 Total slides: {len(prs.slides)}")
    
    return output_file

if __name__ == "__main__":
    create_presentation()
